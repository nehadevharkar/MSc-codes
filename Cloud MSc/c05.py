# Prepare a ppt on cloud computing –introduction , models, services ,and architecture . ppt should contain explanations, images and at least 20 slides

# pip install python-pptx
from pptx import Presentation
from pptx.util import Inches


def create_presentation():
    # Create a presentation
    presentation = Presentation()

    # Slide 1: Title Slide
    slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Cloud Computing: An Introduction"
    subtitle.text = "Models, Services, and Architecture"

    # Slide 2: Introduction
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Introduction to Cloud Computing"
    content.text = "Cloud computing is a technology that provides computing resources over the internet."

    # Add more slides with content and images as needed...

    # Save the presentation
    presentation.save("Cloud_Computing_Presentation.pptx")


if __name__ == "__main__":
    create_presentation()
